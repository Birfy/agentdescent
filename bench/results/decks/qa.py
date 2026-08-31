"""Geometry QA, standing in for a render: LibreOffice cannot convert any pptx here."""
from pptx import Presentation
from pptx.util import Emu
import sys
EMU=914400.0
pr=Presentation('algotune-eight-tasks.pptx')
W=pr.slide_width/EMU; H=pr.slide_height/EMU
# rough width per char at 1pt, Calibri/Cambria ~0.50 em average for mixed case
def est_lines(text, w_in, pt):
    if not text: return 0
    charw = pt*0.50/72.0
    per = max(int(w_in/charw), 1)
    n=0
    for para in text.split('\n'):
        n += max(1, -(-len(para)//per))
    return n
issues=[]
for si,s in enumerate(pr.slides,1):
    boxes=[]
    for sh in s.shapes:
        try: x,y,w,h = sh.left/EMU, sh.top/EMU, sh.width/EMU, sh.height/EMU
        except TypeError: continue
        name=sh.shape_type
        # bounds
        if x < -0.01 or y < -0.01 or x+w > W+0.01 or y+h > H+0.01:
            issues.append(f'slide {si}: OUT OF BOUNDS {name} at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}')
        full = w > W*0.97 and h > H*0.97
        if not sh.has_text_frame: 
            if not full: boxes.append((x,y,w,h,'shape',''))
            continue
        txt = sh.text_frame.text
        if not txt.strip():
            if not full: boxes.append((x,y,w,h,'shape',''))
            continue
        pts=[r.font.size.pt for p in sh.text_frame.paragraphs for r in p.runs if r.font.size]
        pt=max(pts) if pts else 12
        lines=est_lines(txt,w,pt)
        need=lines*pt*1.30/72.0
        if need > h*1.18:
            issues.append(f'slide {si}: TEXT MAY OVERFLOW ~{need:.2f}in into {h:.2f}in box '
                          f'({lines} est lines @ {pt}pt, w={w:.2f}) :: {txt[:58]!r}')
        boxes.append((x,y,w,h,'text',txt[:34]))
    # overlaps between text boxes
    for i in range(len(boxes)):
        for j in range(i+1,len(boxes)):
            a,b=boxes[i],boxes[j]
            if a[4]!='text' or b[4]!='text': continue
            ox=min(a[0]+a[2],b[0]+b[2])-max(a[0],b[0])
            oy=min(a[1]+a[3],b[1]+b[3])-max(a[1],b[1])
            if ox>0.06 and oy>0.06:
                issues.append(f'slide {si}: TEXT OVERLAP {ox:.2f}x{oy:.2f}in :: {a[5]!r} / {b[5]!r}')
print(f'{len(pr.slides)} slides, {W:.2f}x{H:.2f}in')
if not issues: print('no geometry issues found')
for i in issues: print(' ', i)
